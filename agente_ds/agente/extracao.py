# -*- coding: utf-8 -*-
"""
Extrator Documental — Agente DS
Extrai texto de PDFs, DOCX, XLSX, PPTX para análise pelo agente.
Baseado em: 07_AGENTE_1_QUALIFICACAO_DS.md §3 (Sub-Agente 1)

Parsers PDF (ordem de prioridade):
  1. Docling (IBM) — IA para layout, tabelas, hierarquia → Markdown
  2. PyPDF2 — fallback simples se Docling indisponível
"""
import logging
import os
import zipfile
import tempfile
from pathlib import Path
from typing import Dict, List, Optional
import json

log = logging.getLogger("extracao")


class ExtratorDocumental:
    """Extrai e cataloga conteúdo de documentos RFP."""

    EXTENSOES_SUPORTADAS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".csv", ".eml", ".msg", ".html", ".htm"}
    MAX_FILE_SIZE_MB = 100

    def __init__(self):
        self._parsers = {}
        self._docling_available = False
        self._registrar_parsers()

    def _registrar_parsers(self):
        """Registra parsers disponíveis por extensão."""
        # PDF: tentar Docling primeiro (melhor), fallback PyPDF2
        try:
            from docling.document_converter import DocumentConverter
            self._parsers[".pdf"] = self._extrair_pdf_docling
            self._docling_available = True
            log.info("Docling disponível — PDF parser: Docling (IA)")
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                self._parsers[".pdf"] = self._extrair_pdf
                log.info("Docling indisponível — PDF parser: PyPDF2 (fallback)")
            except ImportError:
                pass

        try:
            import docx
            self._parsers[".docx"] = self._extrair_docx
        except ImportError:
            pass

        try:
            import openpyxl
            self._parsers[".xlsx"] = self._extrair_xlsx
        except ImportError:
            pass

        try:
            from pptx import Presentation
            self._parsers[".pptx"] = self._extrair_pptx
        except ImportError:
            pass

        self._parsers[".txt"] = self._extrair_texto
        self._parsers[".md"] = self._extrair_texto
        self._parsers[".csv"] = self._extrair_texto
        self._parsers[".eml"] = self._extrair_eml
        self._parsers[".msg"] = self._extrair_eml
        self._parsers[".html"] = self._extrair_html
        self._parsers[".htm"] = self._extrair_html

    def extrair(self, caminho: str) -> Dict:
        """
        Ponto de entrada principal.
        Aceita: arquivo individual ou ZIP com múltiplos documentos.
        Retorna: ExtractionResult (JSON-serializable dict)
        """
        caminho = Path(caminho)

        if not caminho.exists():
            return self._erro(f"Arquivo não encontrado: {caminho}")

        if caminho.suffix.lower() == ".zip":
            return self._extrair_zip(caminho)
        else:
            return self._extrair_arquivo(caminho)

    def _extrair_zip(self, zip_path: Path) -> Dict:
        """Extrai todos os documentos de um ZIP."""
        documentos = []
        with tempfile.TemporaryDirectory() as tmpdir:
            with zipfile.ZipFile(zip_path, "r") as zf:
                zf.extractall(tmpdir)

            for root, _, files in os.walk(tmpdir):
                for fname in files:
                    fpath = Path(root) / fname
                    ext = fpath.suffix.lower()

                    # Ignorar arquivos ocultos e de sistema
                    if fname.startswith(".") or fname.startswith("~"):
                        continue

                    if ext in self.EXTENSOES_SUPORTADAS:
                        doc = self._extrair_arquivo(fpath)
                        if doc.get("status") == "ok":
                            documentos.append(doc)

        return {
            "status": "ok",
            "fonte": str(zip_path.name),
            "total_documentos": len(documentos),
            "documentos": documentos,
            "texto_consolidado": "\n\n---\n\n".join(
                d.get("texto", "") for d in documentos if d.get("texto")
            ),
        }

    def _extrair_arquivo(self, fpath: Path) -> Dict:
        """Extrai texto de um arquivo individual."""
        ext = fpath.suffix.lower()

        if ext not in self.EXTENSOES_SUPORTADAS:
            return self._erro(f"Extensão não suportada: {ext}")

        # Verificar tamanho
        size_mb = fpath.stat().st_size / (1024 * 1024)
        if size_mb > self.MAX_FILE_SIZE_MB:
            return self._erro(f"Arquivo muito grande: {size_mb:.1f}MB (máx: {self.MAX_FILE_SIZE_MB}MB)")

        parser = self._parsers.get(ext)
        if not parser:
            return self._erro(f"Parser não disponível para {ext}. Instale a dependência.")

        try:
            texto = parser(fpath)
            return {
                "status": "ok",
                "arquivo": fpath.name,
                "extensao": ext,
                "tamanho_mb": round(size_mb, 2),
                "caracteres": len(texto),
                "texto": texto,
            }
        except Exception as e:
            return self._erro(f"Erro ao extrair {fpath.name}: {str(e)}")

    # ─── Parsers por tipo ───

    def _extrair_pdf_docling(self, fpath: Path) -> str:
        """Extrai PDF via Docling (IA) — tabelas e estrutura superiores."""
        from docling.document_converter import DocumentConverter
        try:
            converter = DocumentConverter()
            result = converter.convert(str(fpath))
            markdown = result.document.export_to_markdown()
            if markdown and len(markdown.strip()) > 50:
                log.info(f"Docling extraiu {len(markdown):,} chars de {fpath.name}")
                return markdown
            else:
                log.warning(f"Docling retornou pouco texto, tentando PyPDF2...")
                return self._extrair_pdf(fpath)
        except Exception as e:
            log.warning(f"Docling falhou ({e}), usando PyPDF2 como fallback")
            return self._extrair_pdf(fpath)

    def _extrair_pdf(self, fpath: Path) -> str:
        """Extrai PDF via PyPDF2 (fallback simples)."""
        from PyPDF2 import PdfReader
        reader = PdfReader(str(fpath))
        paginas = []
        for i, page in enumerate(reader.pages):
            texto = page.extract_text() or ""
            if texto.strip():
                paginas.append(f"[Página {i+1}]\n{texto}")
        return "\n\n".join(paginas)

    def _extrair_docx(self, fpath: Path) -> str:
        import docx
        doc = docx.Document(str(fpath))
        paragrafos = [p.text for p in doc.paragraphs if p.text.strip()]

        # Extrair tabelas também
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragrafos.append(" | ".join(cells))

        return "\n".join(paragrafos)

    def _extrair_xlsx(self, fpath: Path) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(str(fpath), data_only=True)
        linhas = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            linhas.append(f"[Aba: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    linhas.append(" | ".join(cells))
        return "\n".join(linhas)

    def _extrair_pptx(self, fpath: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(fpath))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides_text.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides_text)

    def _extrair_texto(self, fpath: Path) -> str:
        encodings = ["utf-8", "latin-1", "cp1252"]
        for enc in encodings:
            try:
                return fpath.read_text(encoding=enc)
            except (UnicodeDecodeError, UnicodeError):
                continue
        return fpath.read_text(encoding="utf-8", errors="replace")

    def _extrair_eml(self, fpath: Path) -> str:
        """Extrai texto de arquivos .eml (e-mails)."""
        import email
        from email import policy

        raw = fpath.read_bytes()
        msg = email.message_from_bytes(raw, policy=policy.default)

        partes = []
        partes.append(f"De: {msg.get('From', '')}")
        partes.append(f"Para: {msg.get('To', '')}")
        partes.append(f"Assunto: {msg.get('Subject', '')}")
        partes.append(f"Data: {msg.get('Date', '')}")
        partes.append("---")

        # Extrair corpo
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    body = part.get_content()
                    if body:
                        partes.append(str(body))
                elif content_type == "text/html":
                    body = part.get_content()
                    if body:
                        partes.append(self._strip_html(str(body)))
        else:
            body = msg.get_content()
            if body:
                partes.append(str(body))

        # Listar anexos
        anexos = []
        for part in msg.walk():
            filename = part.get_filename()
            if filename:
                anexos.append(filename)
        if anexos:
            partes.append(f"\nAnexos: {', '.join(anexos)}")

        return "\n".join(partes)

    def _extrair_html(self, fpath: Path) -> str:
        """Extrai texto de arquivos HTML."""
        raw = self._extrair_texto(fpath)
        return self._strip_html(raw)

    def _strip_html(self, html: str) -> str:
        """Remove tags HTML e retorna texto limpo."""
        import re
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _erro(self, msg: str) -> Dict:
        return {"status": "erro", "mensagem": msg}
